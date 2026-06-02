"""真实产物链 skills 骨架（M3+ / M5）。

提供三个 skill，与 :mod:`agent_hub.pilot.skills.fake` 中的
``fake_slidespec_generate`` / ``fake_pptx_render`` / ``fake_drive_upload``
形成 1:1 替代关系：

- ``real.slide.generate_spec`` —— 基于 brief 文本确定性地生成 SlideSpec
  JSON（章节、要点、layout）。
- ``real.slide.render_pptx`` —— 调用 ``python-pptx`` 渲染为 PPTX 字节流。
  ``python-pptx`` 已固化为主依赖；若运行环境异常缺失，默认严格抛错，
  仅当显式设置环境变量 ``PILOT_ALLOW_PPTX_STUB=true`` 时才降级为 stub
  字节流（artifact metadata 会带 ``is_stub=True`` 明确标注）。
- ``real.drive.upload_share`` —— 通过 :class:`FeishuClientProtocol`
  上传 PPTX 到飞书云空间并返回 share_url。

设计上保持与 fake 对应版本一致的 ``output_artifact_ref`` 与
``input_artifacts`` 协议，便于通过 feature flag 平滑切换。
"""

from __future__ import annotations

import json
import os
import re
import zipfile
from datetime import UTC, datetime
from io import BytesIO
from typing import Any
from xml.sax.saxutils import escape as xml_escape

import structlog

from agent_hub.connectors.feishu.client import FeishuApiError, FeishuClientProtocol
from agent_hub.pilot.domain.enums import ArtifactType
from agent_hub.pilot.skills.internal import ArtifactContentReader
from agent_hub.pilot.skills.registry import (
    SkillInvocation,
    SkillRegistry,
    SkillResult,
    SkillSpec,
)

REAL_CHAIN_SCOPES = ["pilot", "pilot.real_chain"]
logger = structlog.get_logger(__name__)


# ── real.slide.generate_spec ───────────────────────


def _build_slide_spec(brief_text: str, title: str) -> dict[str, Any]:
    """Build a practical SlideSpec from markdown brief content.

    This stays deterministic so the renderer remains testable; the content
    quality comes from the upstream brief generator when an LLM is configured.
    """
    deck_title = _resolve_title(title, brief_text, default="Generated Deck")
    sections = _parse_markdown_sections(brief_text, fallback_title=deck_title)
    sections = [
        section
        for section in sections
        if section["title"] != deck_title or section["bullets"]
    ]
    if not sections:
        sections = [{
            "title": "核心内容",
            "bullets": ["暂无可用内容，请补充需求背景、目标和交付约束。"],
        }]

    subtitle = _extract_subtitle(brief_text) or "AI agent generated working deck"
    slides: list[dict[str, Any]] = [
        {
            "layout": "cover",
            "title": deck_title,
            "subtitle": subtitle,
            "kicker": "Generated Deck",
        },
    ]

    if len(sections) >= 3:
        slides.append(
            {
                "layout": "agenda",
                "title": "目录",
                "items": [section["title"] for section in sections[:6]],
            },
        )

    for index, section in enumerate(sections[:10], start=1):
        bullets = _clean_bullets(section["bullets"], limit=6)
        layout = "two_column" if len(bullets) >= 5 else "bullets"
        slides.append(
            {
                "layout": layout,
                "kicker": f"{index:02d}",
                "title": section["title"],
                "bullets": bullets or ["待补充具体要点。"],
                "speaker_notes": f"Explain {section['title']} with concrete context.",
            },
        )

    slides.append(
        {
            "layout": "closing",
            "title": "下一步",
            "bullets": _derive_next_steps(sections),
        },
    )
    return {
        "title": deck_title,
        "slide_count": len(slides),
        "theme": "agent_loop_v1",
        "slides": slides,
    }


def _resolve_title(title: str, markdown: str, *, default: str) -> str:
    explicit = (title or "").strip()
    if explicit and explicit.lower() not in {"untitled deck", "deck", "document"}:
        return explicit
    inferred = _first_markdown_heading(markdown)
    return inferred or explicit or default


def _first_markdown_heading(markdown: str) -> str:
    for line in (markdown or "").splitlines():
        stripped = line.strip()
        if stripped.startswith("# "):
            return _clean_inline_markdown(stripped[2:])[:80]
    return ""


def _extract_subtitle(markdown: str) -> str:
    for line in (markdown or "").splitlines():
        stripped = line.strip()
        if stripped.startswith(">"):
            return _clean_inline_markdown(stripped.lstrip("> ").strip())[:120]
    return ""


def _parse_markdown_sections(
    markdown: str,
    *,
    fallback_title: str,
) -> list[dict[str, Any]]:
    sections: list[dict[str, Any]] = []
    current: dict[str, Any] | None = None
    for raw_line in (markdown or "").splitlines():
        line = raw_line.strip()
        if not line:
            continue
        heading = _heading_text(line)
        if heading:
            if current:
                sections.append(current)
            current = {"title": heading, "bullets": []}
            continue
        if current is None:
            current = {"title": fallback_title, "bullets": []}
        current["bullets"].append(_clean_inline_markdown(_strip_list_marker(line)))
    if current:
        sections.append(current)
    return sections


def _heading_text(line: str) -> str:
    match = re.match(r"^(#{1,3})\s+(.+)$", line)
    if not match:
        return ""
    return _clean_inline_markdown(match.group(2))[:80]


def _strip_list_marker(line: str) -> str:
    line = re.sub(r"^\s*[-*+]\s+", "", line)
    line = re.sub(r"^\s*\d+[.)]\s+", "", line)
    return line


def _clean_inline_markdown(text: str) -> str:
    cleaned = re.sub(r"\*\*(.*?)\*\*", r"\1", text)
    cleaned = re.sub(r"`([^`]+)`", r"\1", cleaned)
    cleaned = re.sub(r"\[(.*?)\]\([^)]*\)", r"\1", cleaned)
    return cleaned.strip()


def _clean_bullets(values: list[Any], *, limit: int) -> list[str]:
    bullets: list[str] = []
    for value in values:
        text = _clean_inline_markdown(str(value)).strip()
        if not text or text.startswith("<!--"):
            continue
        if len(text) > 120:
            text = text[:117].rstrip() + "..."
        bullets.append(text)
        if len(bullets) >= limit:
            break
    return bullets


def _derive_next_steps(sections: list[dict[str, Any]]) -> list[str]:
    for section in sections:
        if "下一" in section["title"] or "行动" in section["title"]:
            bullets = _clean_bullets(section["bullets"], limit=4)
            if bullets:
                return bullets
    return ["确认叙事主线与关键数据", "补充案例、截图或引用来源", "按审批结果交付并同步飞书链接"]


def _make_real_slidespec_skill() -> tuple[SkillSpec, Any]:
    async def _gen(inv: SkillInvocation) -> SkillResult:
        params = inv.params
        title = str(params.get("title", "") or "")
        upstream = params.get("input_artifacts", {}) or {}
        brief_artifact = upstream.get("brief")
        brief_text = ""
        if isinstance(brief_artifact, dict):
            content = brief_artifact.get("content")
            if isinstance(content, str):
                brief_text = content
            elif isinstance(content, dict | list):
                brief_text = json.dumps(content, ensure_ascii=False)

        spec = _build_slide_spec(brief_text, title)

        if inv.dry_run:
            return SkillResult(
                skill_name=inv.skill_name,
                success=True,
                output={"slide_count": spec["slide_count"], "would_generate": True},
            )

        return SkillResult(
            skill_name=inv.skill_name,
            success=True,
            output={"slide_count": spec["slide_count"]},
            artifact_payload={
                "type": ArtifactType.SLIDE_SPEC_JSON.value,
                "title": f"{spec['title']} SlideSpec",
                "mime_type": "application/json",
                "content": spec,
                "metadata": {
                    "slide_count": spec["slide_count"],
                    "theme": spec.get("theme", ""),
                },
            },
        )

    spec = SkillSpec(
        skill_name="real.slide.generate_spec",
        description="基于 brief 内容确定性生成 SlideSpec JSON。",
        side_effect="read",
        requires_approval=False,
        supports_dry_run=True,
        scopes=REAL_CHAIN_SCOPES,
        timeout_ms=10_000,
    )
    return spec, _gen


# ── real.slide.render_pptx ─────────────────────────


def _render_pptx_bytes(slide_spec: dict[str, Any]) -> tuple[bytes, bool]:
    """Render a styled, openable PPTX; return ``(bytes, is_stub)``."""
    try:
        from pptx import Presentation  # type: ignore[import-not-found]
        from pptx.dml.color import RGBColor  # type: ignore[import-not-found]
        from pptx.enum.shapes import MSO_SHAPE  # type: ignore[import-not-found]
        from pptx.util import Inches, Pt  # type: ignore[import-not-found]
    except ImportError as exc:
        allow_stub = os.environ.get("PILOT_ALLOW_PPTX_STUB", "").lower() in (
            "1",
            "true",
            "yes",
        )
        if not allow_stub:
            raise RuntimeError(
                "real.slide.render_pptx needs python-pptx; install it or set "
                "PILOT_ALLOW_PPTX_STUB=true for explicit demo fallback.",
            ) from exc
        logger.warning(
            "real.slide.render_pptx.fallback_to_stub",
            reason="python-pptx not installed",
        )
        stub = ("PPTX_STUB::" + json.dumps(slide_spec, ensure_ascii=False)).encode(
            "utf-8",
        )
        return stub, True

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    bg = RGBColor(248, 250, 252)
    ink = RGBColor(15, 23, 42)
    muted = RGBColor(71, 85, 105)
    accent = RGBColor(14, 165, 163)
    accent_dark = RGBColor(15, 118, 110)
    warm = RGBColor(245, 158, 11)

    def add_text(slide, x, y, w, h, text, size, *, bold=False, color=None, align=None):
        shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        if align is not None:
            p.alignment = align
        run = p.add_run()
        run.text = str(text or "")
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color or ink
        return shape

    def add_rule(slide, x, y, w, h, color):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(w),
            Inches(h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def add_card(slide, x, y, w, h, color=None):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(w),
            Inches(h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color or RGBColor(255, 255, 255)
        shape.line.color.rgb = RGBColor(226, 232, 240)
        return shape

    def add_bullets(slide, x, y, w, h, bullets, *, size=21):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        for index, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.name = "Microsoft YaHei"
            p.font.size = Pt(size)
            p.font.color.rgb = ink
            p.space_after = Pt(9)
        return box

    def paint_background(slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = bg

    slides = list(slide_spec.get("slides", []) or [])
    if not slides:
        slides = [{"layout": "cover", "title": slide_spec.get("title", "Deck")}]

    for index, item in enumerate(slides):
        slide = prs.slides.add_slide(blank_layout)
        paint_background(slide)
        layout = str(item.get("layout") or "bullets")
        title = str(item.get("title") or slide_spec.get("title") or "Deck")
        bullets = [str(b) for b in (item.get("bullets") or item.get("items") or [])]

        add_rule(slide, 0, 0, 0.18, 7.5, accent)
        add_text(slide, 0.55, 0.28, 2.5, 0.35, f"{index + 1:02d}", 12, color=muted)

        if layout in {"cover", "title"}:
            add_rule(slide, 0.75, 1.15, 0.08, 3.0, accent)
            add_text(slide, 1.05, 1.45, 10.7, 1.35, title, 42, bold=True, color=ink)
            add_text(
                slide,
                1.07,
                2.9,
                9.7,
                0.85,
                item.get("subtitle") or "Generated by Pilot agent loop",
                19,
                color=muted,
            )
            add_card(slide, 1.05, 5.7, 3.6, 0.55, RGBColor(236, 253, 245))
            add_text(slide, 1.25, 5.84, 3.2, 0.25, item.get("kicker") or "Working draft", 11, color=accent_dark)
            continue

        if layout == "agenda":
            add_text(slide, 0.85, 0.75, 8.5, 0.7, title, 30, bold=True)
            for i, label in enumerate(bullets[:6], start=1):
                y = 1.65 + (i - 1) * 0.72
                add_card(slide, 0.95, y, 10.8, 0.5)
                add_text(slide, 1.15, y + 0.11, 0.5, 0.25, f"{i}", 14, bold=True, color=warm)
                add_text(slide, 1.8, y + 0.09, 9.3, 0.28, label, 16, color=ink)
            continue

        add_text(slide, 0.85, 0.55, 0.7, 0.4, item.get("kicker") or "", 15, bold=True, color=accent_dark)
        add_text(slide, 1.55, 0.5, 10.7, 0.65, title, 28, bold=True)
        add_rule(slide, 1.55, 1.28, 3.0, 0.04, accent)

        if layout == "two_column":
            midpoint = (len(bullets) + 1) // 2
            left, right = bullets[:midpoint], bullets[midpoint:]
            add_card(slide, 0.9, 1.75, 5.65, 4.6)
            add_card(slide, 6.75, 1.75, 5.65, 4.6)
            add_bullets(slide, 1.25, 2.05, 4.95, 3.95, left, size=19)
            add_bullets(slide, 7.1, 2.05, 4.95, 3.95, right, size=19)
        elif layout == "closing":
            for i, bullet in enumerate(bullets[:4], start=1):
                y = 1.8 + (i - 1) * 0.95
                add_card(slide, 1.05, y, 10.7, 0.72, RGBColor(255, 251, 235))
                add_text(slide, 1.32, y + 0.19, 0.45, 0.25, f"{i}", 15, bold=True, color=warm)
                add_text(slide, 1.95, y + 0.17, 9.0, 0.32, bullet, 18)
        else:
            add_card(slide, 0.95, 1.7, 11.3, 4.8)
            add_bullets(slide, 1.35, 2.05, 10.4, 4.1, bullets[:6], size=21)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue(), False


def _make_real_pptx_render_skill() -> tuple[SkillSpec, Any]:
    async def _render(inv: SkillInvocation) -> SkillResult:
        params = inv.params
        upstream = params.get("input_artifacts", {}) or {}
        spec_artifact = upstream.get("spec") or upstream.get("slide_spec")
        slide_spec: dict[str, Any] | None = None
        if isinstance(spec_artifact, dict):
            content = spec_artifact.get("content")
            if isinstance(content, dict):
                slide_spec = content
            elif isinstance(content, str):
                try:
                    parsed = json.loads(content)
                    if isinstance(parsed, dict):
                        slide_spec = parsed
                except json.JSONDecodeError:
                    pass

        if slide_spec is None:
            return SkillResult(
                skill_name=inv.skill_name,
                success=False,
                error="real.slide.render_pptx 需要 input_artifacts.spec",
            )

        title = str(params.get("title", "") or slide_spec.get("title", "Deck"))

        if inv.dry_run:
            return SkillResult(
                skill_name=inv.skill_name,
                success=True,
                output={
                    "slide_count": slide_spec.get("slide_count", 0),
                    "would_render": True,
                },
            )

        pptx_bytes, is_stub = _render_pptx_bytes(slide_spec)

        return SkillResult(
            skill_name=inv.skill_name,
            success=True,
            output={"byte_size": len(pptx_bytes), "is_stub": is_stub},
            artifact_payload={
                "type": ArtifactType.PPTX.value,
                "title": f"{title}.pptx",
                "mime_type": (
                    "application/vnd.openxmlformats-officedocument"
                    ".presentationml.presentation"
                ),
                "content": pptx_bytes,
                "metadata": {
                    "slide_count": slide_spec.get("slide_count", 0),
                    "theme": slide_spec.get("theme", "agent_loop_v1"),
                    "renderer": "pptx-stub" if is_stub else "python-pptx",
                    "is_stub": is_stub,
                },
            },
        )

    spec = SkillSpec(
        skill_name="real.slide.render_pptx",
        description="将 SlideSpec 渲染为 PPTX（python-pptx；缺失时降级为 stub）。",
        side_effect="write",
        requires_approval=True,
        supports_dry_run=True,
        scopes=REAL_CHAIN_SCOPES,
        idempotency_key_fields=["title"],
        timeout_ms=30_000,
    )
    return spec, _render


# ── real.doc.render_docx ───────────────────────────


def _render_docx_bytes(markdown_text: str, title: str) -> tuple[bytes, str]:
    doc_title = _resolve_title(title, markdown_text, default="Document")
    blocks = _markdown_to_doc_blocks(markdown_text, doc_title)
    created = datetime.now(UTC).replace(microsecond=0).isoformat().replace("+00:00", "Z")

    buf = BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("[Content_Types].xml", _docx_content_types_xml())
        zf.writestr("_rels/.rels", _docx_root_rels_xml())
        zf.writestr("docProps/core.xml", _docx_core_xml(doc_title, created))
        zf.writestr("docProps/app.xml", _docx_app_xml())
        zf.writestr("word/document.xml", _docx_document_xml(blocks))
        zf.writestr("word/styles.xml", _docx_styles_xml())
    return buf.getvalue(), doc_title


def _markdown_to_doc_blocks(
    markdown_text: str,
    title: str,
) -> list[tuple[str, str]]:
    blocks: list[tuple[str, str]] = []
    for raw_line in (markdown_text or "").splitlines():
        line = raw_line.strip()
        if not line:
            continue
        heading = _heading_text(line)
        if heading:
            level = len(line) - len(line.lstrip("#"))
            style = "Title" if level == 1 else f"Heading{min(level - 1, 2)}"
            blocks.append((style, heading))
            continue
        if re.match(r"^\s*([-*+]|\d+[.)])\s+", raw_line):
            blocks.append(("ListBullet", _clean_inline_markdown(_strip_list_marker(line))))
            continue
        if line.startswith(">"):
            blocks.append(("Quote", _clean_inline_markdown(line.lstrip("> ").strip())))
            continue
        blocks.append(("Normal", _clean_inline_markdown(line)))

    if not blocks or blocks[0] != ("Title", title):
        blocks.insert(0, ("Title", title))
    return blocks


def _docx_document_xml(blocks: list[tuple[str, str]]) -> str:
    body = "\n".join(_docx_paragraph(style, text) for style, text in blocks if text)
    return f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
  <w:body>
    {body}
    <w:sectPr>
      <w:pgSz w:w="11906" w:h="16838"/>
      <w:pgMar w:top="1440" w:right="1440" w:bottom="1440" w:left="1440"/>
    </w:sectPr>
  </w:body>
</w:document>"""


def _docx_paragraph(style: str, text: str) -> str:
    escaped = xml_escape(text)
    return f"""<w:p>
  <w:pPr><w:pStyle w:val="{style}"/></w:pPr>
  <w:r><w:t xml:space="preserve">{escaped}</w:t></w:r>
</w:p>"""


def _docx_content_types_xml() -> str:
    return """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/docProps/core.xml" ContentType="application/vnd.openxmlformats-package.core-properties+xml"/>
  <Override PartName="/docProps/app.xml" ContentType="application/vnd.openxmlformats-officedocument.extended-properties+xml"/>
  <Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
  <Override PartName="/word/styles.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.styles+xml"/>
</Types>"""


def _docx_root_rels_xml() -> str:
    return """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
  <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/package/2006/relationships/metadata/core-properties" Target="docProps/core.xml"/>
  <Relationship Id="rId3" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/extended-properties" Target="docProps/app.xml"/>
</Relationships>"""


def _docx_core_xml(title: str, created: str) -> str:
    escaped_title = xml_escape(title)
    return f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<cp:coreProperties xmlns:cp="http://schemas.openxmlformats.org/package/2006/metadata/core-properties" xmlns:dc="http://purl.org/dc/elements/1.1/" xmlns:dcterms="http://purl.org/dc/terms/" xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance">
  <dc:title>{escaped_title}</dc:title>
  <dc:creator>Pilot Agent</dc:creator>
  <cp:lastModifiedBy>Pilot Agent</cp:lastModifiedBy>
  <dcterms:created xsi:type="dcterms:W3CDTF">{created}</dcterms:created>
  <dcterms:modified xsi:type="dcterms:W3CDTF">{created}</dcterms:modified>
</cp:coreProperties>"""


def _docx_app_xml() -> str:
    return """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Properties xmlns="http://schemas.openxmlformats.org/officeDocument/2006/extended-properties" xmlns:vt="http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes">
  <Application>Pilot Agent</Application>
</Properties>"""


def _docx_styles_xml() -> str:
    return """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:styles xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
  <w:docDefaults>
    <w:rPrDefault><w:rPr><w:rFonts w:ascii="Microsoft YaHei" w:eastAsia="Microsoft YaHei"/><w:sz w:val="22"/></w:rPr></w:rPrDefault>
  </w:docDefaults>
  <w:style w:type="paragraph" w:styleId="Normal"><w:name w:val="Normal"/></w:style>
  <w:style w:type="paragraph" w:styleId="Title"><w:name w:val="Title"/><w:pPr><w:spacing w:after="240"/></w:pPr><w:rPr><w:b/><w:sz w:val="44"/></w:rPr></w:style>
  <w:style w:type="paragraph" w:styleId="Heading1"><w:name w:val="heading 1"/><w:pPr><w:spacing w:before="240" w:after="120"/></w:pPr><w:rPr><w:b/><w:color w:val="0F766E"/><w:sz w:val="32"/></w:rPr></w:style>
  <w:style w:type="paragraph" w:styleId="Heading2"><w:name w:val="heading 2"/><w:pPr><w:spacing w:before="180" w:after="80"/></w:pPr><w:rPr><w:b/><w:sz w:val="28"/></w:rPr></w:style>
  <w:style w:type="paragraph" w:styleId="ListBullet"><w:name w:val="List Bullet"/><w:pPr><w:ind w:left="720" w:hanging="360"/></w:pPr></w:style>
  <w:style w:type="paragraph" w:styleId="Quote"><w:name w:val="Quote"/><w:pPr><w:ind w:left="360"/><w:spacing w:after="120"/></w:pPr><w:rPr><w:i/><w:color w:val="475569"/></w:rPr></w:style>
</w:styles>"""


def _make_real_docx_render_skill() -> tuple[SkillSpec, Any]:
    async def _render(inv: SkillInvocation) -> SkillResult:
        params = inv.params
        upstream = params.get("input_artifacts", {}) or {}
        brief_artifact = upstream.get("brief") or upstream.get("doc") or upstream.get("file")
        markdown_text = ""
        if isinstance(brief_artifact, dict):
            content = brief_artifact.get("content")
            if isinstance(content, str):
                markdown_text = content
            elif isinstance(content, bytes):
                markdown_text = content.decode("utf-8", errors="replace")
            elif isinstance(content, dict | list):
                markdown_text = json.dumps(content, ensure_ascii=False, indent=2)

        if not markdown_text.strip():
            return SkillResult(
                skill_name=inv.skill_name,
                success=False,
                error="real.doc.render_docx 需要 input_artifacts.brief / doc / file",
            )

        title = str(params.get("title", "") or "")
        doc_title = _resolve_title(title, markdown_text, default="Document")
        if inv.dry_run:
            return SkillResult(
                skill_name=inv.skill_name,
                success=True,
                output={"title": doc_title, "would_render": True},
            )

        docx_bytes, doc_title = _render_docx_bytes(markdown_text, doc_title)
        return SkillResult(
            skill_name=inv.skill_name,
            success=True,
            output={"byte_size": len(docx_bytes), "title": doc_title},
            artifact_payload={
                "type": ArtifactType.DOCX.value,
                "title": f"{doc_title}.docx",
                "mime_type": (
                    "application/vnd.openxmlformats-officedocument"
                    ".wordprocessingml.document"
                ),
                "content": docx_bytes,
                "metadata": {
                    "renderer": "docx-ooxml",
                    "char_count": len(markdown_text),
                },
            },
        )

    spec = SkillSpec(
        skill_name="real.doc.render_docx",
        description="将 Markdown Brief 渲染为可打开的 Word DOCX 文件。",
        side_effect="write",
        requires_approval=True,
        supports_dry_run=True,
        scopes=REAL_CHAIN_SCOPES,
        idempotency_key_fields=["title"],
        timeout_ms=30_000,
    )
    return spec, _render


# ── real.drive.upload_share ────────────────────────


def _make_real_upload_share_skill(
    client: FeishuClientProtocol,
    *,
    artifact_reader: ArtifactContentReader,
    default_folder_token: str = "",
    admin_open_id: str = "",
    drive_url_template: str = "https://www.feishu.cn/file/{file_token}",
) -> tuple[SkillSpec, Any]:
    async def _upload(inv: SkillInvocation) -> SkillResult:
        params = inv.params
        file_name = str(params.get("file_name", "") or "deck.pptx")
        folder_token = str(params.get("folder_token", "") or default_folder_token) or None
        recipient_open_ids = _collect_open_ids(
            params.get("member_open_id"),
            params.get("share_recipient_open_id"),
            params.get("recipient_open_id"),
            params.get("member_open_ids"),
            params.get("share_recipient_open_ids"),
            params.get("recipient_open_ids"),
            params.get("shared_open_ids"),
        )

        upstream = params.get("input_artifacts", {}) or {}
        # 兼容 PPTX / 任意文件 / Brief Markdown / Doc 上游，让该 skill 可以复用为
        # DOC_ONLY plan 的交付环节。
        source_artifact = (
            upstream.get("pptx")
            or upstream.get("docx")
            or upstream.get("file")
            or upstream.get("brief")
            or upstream.get("doc")
        )

        content_bytes: bytes | None = None
        if isinstance(source_artifact, dict):
            raw = source_artifact.get("content")
            if isinstance(raw, bytes):
                content_bytes = raw
            elif isinstance(raw, str):
                content_bytes = raw.encode("utf-8")

        if content_bytes is None and "artifact_id" in params:
            raw = await artifact_reader(str(params["artifact_id"]))
            if isinstance(raw, bytes):
                content_bytes = raw
            elif isinstance(raw, str):
                content_bytes = raw.encode("utf-8")

        if content_bytes is None:
            return SkillResult(
                skill_name=inv.skill_name,
                success=False,
                error="real.drive.upload_share 需要 input_artifacts.pptx / docx / brief / doc / file",
            )

        if inv.dry_run:
            return SkillResult(
                skill_name=inv.skill_name,
                success=True,
                output={
                    "file_name": file_name,
                    "byte_size": len(content_bytes),
                    "would_upload": True,
                    "share_recipient_open_ids": recipient_open_ids
                    or ([admin_open_id] if admin_open_id else []),
                },
            )

        try:
            uploaded = await client.upload_file(
                file_name=file_name,
                content=content_bytes,
                parent_type="explorer",
                parent_node=folder_token,
            )
        except FeishuApiError as exc:
            return SkillResult(
                skill_name=inv.skill_name,
                success=False,
                error=str(exc),
                output={
                    "file_name": file_name,
                    "byte_size": len(content_bytes),
                    "folder_token": folder_token or "",
                },
                error_details={
                    "phase": "upload_file",
                    "file_name": file_name,
                    "folder_token": folder_token or "",
                    "feishu": exc.to_error_details(),
                },
            )

        share_url = (uploaded.download_url or "").strip()
        # 飞书 Drive upload_all 实际并不返回稳定可访问 URL，需要基于 file_token 兜底拼接。
        if not share_url and uploaded.file_token and drive_url_template:
            try:
                share_url = drive_url_template.format(file_token=uploaded.file_token)
            except (KeyError, IndexError):
                share_url = ""

        # Prefer explicit requester/shared open_ids for per-task delivery; keep
        # admin_open_id only as a backward-compatible fallback.
        # need_notification=True 让飞书主动发一条带可点击链接的"与我共享"通知。
        share_member_open_ids = recipient_open_ids or (
            [admin_open_id] if admin_open_id else []
        )
        shared_open_ids: list[str] = []
        if uploaded.file_token:
            for open_id in share_member_open_ids:
                try:
                    await client.share_file(
                        file_token=uploaded.file_token,
                        member_open_id=open_id,
                        perm="edit",
                        need_notification=True,
                        file_type="file",
                    )
                    shared_open_ids.append(open_id)
                except FeishuApiError as exc:
                    logger.warning(
                        "drive_share_failed",
                        file_token=uploaded.file_token,
                        member_open_id=open_id,
                        error=str(exc),
                    )
                    return SkillResult(
                        skill_name=inv.skill_name,
                        success=False,
                        error=str(exc),
                        output={
                            "file_token": uploaded.file_token,
                            "share_url": share_url,
                            "byte_size": uploaded.size,
                            "folder_token": folder_token or "",
                            "shared_open_ids": shared_open_ids,
                        },
                        error_details={
                            "phase": "share_file",
                            "file_token": uploaded.file_token,
                            "member_open_id": open_id,
                            "shared_open_ids": shared_open_ids,
                            "feishu": exc.to_error_details(),
                        },
                    )

        return SkillResult(
            skill_name=inv.skill_name,
            success=True,
            output={
                "file_token": uploaded.file_token,
                "share_url": share_url,
                "byte_size": uploaded.size,
            },
            artifact_payload={
                "type": ArtifactType.SHARE_LINK.value,
                "title": f"{file_name} share link",
                "uri": share_url,
                "mime_type": "text/uri-list",
                "content": {
                    "file_token": uploaded.file_token,
                    "share_url": share_url,
                    "file_name": file_name,
                    "shared_open_id": share_member_open_ids[0]
                    if share_member_open_ids
                    else "",
                    "shared_open_ids": share_member_open_ids,
                },
                "metadata": {
                    "file_token": uploaded.file_token,
                    "feishu_token": uploaded.file_token,
                    "byte_size": uploaded.size,
                    "folder_token": folder_token or "",
                    "share_url": share_url,
                    "shared_open_id": share_member_open_ids[0]
                    if share_member_open_ids
                    else "",
                    "shared_open_ids": share_member_open_ids,
                },
            },
        )

    spec = SkillSpec(
        skill_name="real.drive.upload_share",
        description="把 PPTX 上传到飞书云空间并返回（暂用 download_url 充当）share_url。",
        side_effect="share",
        requires_approval=True,
        supports_dry_run=True,
        scopes=REAL_CHAIN_SCOPES,
        idempotency_key_fields=["file_name", "folder_token"],
        timeout_ms=45_000,
    )
    return spec, _upload


def _collect_open_ids(*values: object) -> list[str]:
    """Collect comma/list-style open_id inputs while preserving order."""
    result: list[str] = []
    seen: set[str] = set()
    for value in values:
        for item in _iter_open_ids(value):
            if item and item not in seen:
                seen.add(item)
                result.append(item)
    return result


def _iter_open_ids(value: object):  # noqa: ANN202
    if value is None:
        return
    if isinstance(value, str):
        raw = value.strip()
        if not raw:
            return
        if raw.startswith("["):
            try:
                parsed = json.loads(raw)
            except json.JSONDecodeError:
                parsed = None
            if parsed is not None:
                yield from _iter_open_ids(parsed)
                return
        for part in raw.replace("；", ",").replace(";", ",").split(","):
            item = part.strip()
            if item:
                yield item
        return
    if isinstance(value, dict):
        for key in ("open_id", "member_open_id", "user_id"):
            if key in value:
                yield from _iter_open_ids(value[key])
        return
    try:
        iterator = iter(value)
    except TypeError:
        return
    for item in iterator:
        yield from _iter_open_ids(item)


# ── 注册入口 ────────────────────────────────────────


def register_real_chain_skills(
    registry: SkillRegistry,
    *,
    feishu_client: FeishuClientProtocol | None = None,
    artifact_reader: ArtifactContentReader | None = None,
    default_folder_token: str = "",
    admin_open_id: str = "",
    drive_url_template: str = "https://www.feishu.cn/file/{file_token}",
    allow_overwrite: bool = False,
) -> None:
    """注册真实产物链 skills。

    SlideSpec / PPTX 不依赖外部服务；upload_share 仅在传入飞书 client
    时注册。便于在没有飞书凭据的环境下也能跑前两步。
    """
    spec, fn = _make_real_slidespec_skill()
    registry.register(spec, fn, allow_overwrite=allow_overwrite)

    spec, fn = _make_real_pptx_render_skill()
    registry.register(spec, fn, allow_overwrite=allow_overwrite)

    spec, fn = _make_real_docx_render_skill()
    registry.register(spec, fn, allow_overwrite=allow_overwrite)

    if feishu_client is not None and artifact_reader is not None:
        spec, fn = _make_real_upload_share_skill(
            feishu_client,
            artifact_reader=artifact_reader,
            default_folder_token=default_folder_token,
            admin_open_id=admin_open_id,
            drive_url_template=drive_url_template,
        )
        registry.register(spec, fn, allow_overwrite=allow_overwrite)


__all__ = [
    "REAL_CHAIN_SCOPES",
    "register_real_chain_skills",
]
